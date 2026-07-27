"""
Base classes and shared formatting methods for PowerPoint charts.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.util import Cm, Pt

from .references import AXIS_POSITION, LABEL_POSITION, LEGEND_POSITION
from .utils import validate_option


@dataclass
class BasePowerPointChart:
    """
    Base class for PowerPoint chart objects.

    This class centralizes common chart behavior such as slide validation,
    chart positioning, title formatting, legend formatting, axis formatting
    and plot data labels.
    """

    presentation: Presentation | None = None
    slide: int = 0

    position_x: float = 5
    position_y: float = 5
    width: float = 10
    height: float = 10

    title_show: bool = False
    title: str = ""
    font_name: str = "Inter"
    title_font_size: float = 6

    legend_show: bool = False
    legend_font_size: float = 5
    legend_position: str = "bottom"

    def validate_base_inputs(self) -> None:
        """
        Validate common chart inputs.

        Raises
        ------
        ValueError
            If presentation is missing, slide is out of range or legend_position
            is not valid.
        """
        if self.presentation is None:
            raise ValueError("presentation cannot be None.")

        if self.slide < 0:
            raise ValueError("slide must be greater than or equal to zero.")

        if self.slide >= len(self.presentation.slides):
            raise ValueError(
                f"slide {self.slide} is out of range. "
                f"The presentation has {len(self.presentation.slides)} slides."
            )

        validate_option(
            value=self.legend_position,
            options=LEGEND_POSITION,
            parameter_name="legend_position",
        )

    def get_slide(self):
        """
        Return the PowerPoint slide where the chart will be inserted.

        Returns
        -------
        pptx.slide.Slide
            Selected PowerPoint slide.
        """
        self.validate_base_inputs()
        return self.presentation.slides[self.slide]

    def get_chart_position(self):
        """
        Convert chart position and dimensions from centimeters to EMUs.

        Returns
        -------
        tuple
            Tuple containing x, y, width and height as python-pptx length units.
        """
        return (
            Cm(self.position_x),
            Cm(self.position_y),
            Cm(self.width),
            Cm(self.height),
        )

    def apply_title(self, chart: Chart) -> None:
        """
        Apply title visibility and formatting.

        Parameters
        ----------
        chart : pptx.chart.chart.Chart
            Chart object to format.
        """
        chart.has_title = self.title_show

        if not self.title_show:
            return

        chart.chart_title.text_frame.text = self.title
        paragraph = chart.chart_title.text_frame.paragraphs[0]
        paragraph.font.name = self.font_name
        paragraph.font.size = Pt(self.title_font_size)

    def apply_legend(self, chart: Chart) -> None:
        """
        Apply legend visibility, position and font formatting.

        Parameters
        ----------
        chart : pptx.chart.chart.Chart
            Chart object to format.
        """
        chart.has_legend = self.legend_show

        if not self.legend_show:
            return

        chart.legend.position = LEGEND_POSITION[self.legend_position]
        chart.legend.font.size = Pt(self.legend_font_size)
        chart.legend.font.name = self.font_name

    def apply_value_axis(
        self,
        chart: Chart,
        title_show: bool = False,
        title: str = "",
        font_size: float = 6,
        number_format: str = "0.00",
        show_gridlines: bool = False,
        tick_label_position: str | None = None,
        minimum_scale: float | None = None,
        maximum_scale: float | None = None,
    ) -> None:
        """
        Apply formatting to the value axis.

        Parameters
        ----------
        chart : pptx.chart.chart.Chart
            Chart object to format.
        title_show : bool, default False
            Whether to display the axis title.
        title : str, default ""
            Axis title text.
        font_size : float, default 6
            Font size used for title and tick labels.
        number_format : str, default "0.00"
            Axis tick label number format.
        show_gridlines : bool, default False
            Whether to display major gridlines.
        tick_label_position : str, optional
            Tick label position key from AXIS_POSITION.
        minimum_scale : float, optional
            Minimum axis scale.
        maximum_scale : float, optional
            Maximum axis scale.
        """
        axis = chart.value_axis
        axis.has_title = title_show

        if tick_label_position is not None:
            validate_option(
                value=tick_label_position,
                options=AXIS_POSITION,
                parameter_name="tick_label_position",
            )
            axis.tick_label_position = AXIS_POSITION[tick_label_position]

        if title_show:
            axis.axis_title.text_frame.text = title
            paragraph = axis.axis_title.text_frame.paragraphs[0]
            paragraph.font.name = self.font_name
            paragraph.font.size = Pt(font_size)

        axis.tick_labels.font.size = Pt(font_size)
        axis.tick_labels.font.name = self.font_name
        axis.tick_labels.number_format = number_format
        axis.has_major_gridlines = show_gridlines

        if minimum_scale is not None:
            axis.minimum_scale = float(minimum_scale)

        if maximum_scale is not None:
            axis.maximum_scale = float(maximum_scale)

    def apply_category_axis(
        self,
        chart: Chart,
        title_show: bool = False,
        title: str = "",
        font_size: float = 6,
        number_format: str = "0.00",
        tick_label_position: str | None = None,
        minimum_scale: float | None = None,
        maximum_scale: float | None = None,
    ) -> None:
        """
        Apply formatting to the category axis.

        Parameters
        ----------
        chart : pptx.chart.chart.Chart
            Chart object to format.
        title_show : bool, default False
            Whether to display the axis title.
        title : str, default ""
            Axis title text.
        font_size : float, default 6
            Font size used for title and tick labels.
        number_format : str, default "0.00"
            Axis tick label number format.
        tick_label_position : str, optional
            Tick label position key from AXIS_POSITION.
        minimum_scale : float, optional
            Minimum axis scale. Useful for bubble charts.
        maximum_scale : float, optional
            Maximum axis scale. Useful for bubble charts.
        """
        axis = chart.category_axis
        axis.has_title = title_show

        if tick_label_position is not None:
            validate_option(
                value=tick_label_position,
                options=AXIS_POSITION,
                parameter_name="tick_label_position",
            )
            axis.tick_label_position = AXIS_POSITION[tick_label_position]

        if title_show:
            axis.axis_title.text_frame.text = title
            paragraph = axis.axis_title.text_frame.paragraphs[0]
            paragraph.font.name = self.font_name
            paragraph.font.size = Pt(font_size)

        axis.tick_labels.font.size = Pt(font_size)
        axis.tick_labels.font.name = self.font_name
        axis.tick_labels.number_format = number_format

        if minimum_scale is not None:
            axis.minimum_scale = float(minimum_scale)

        if maximum_scale is not None:
            axis.maximum_scale = float(maximum_scale)

    def apply_plot_data_labels(
        self,
        chart: Chart,
        show_labels: bool = False,
        label_position: str = "center",
        number_format: str = "0.0",
        font_size: float = 5,
        show_category_name: bool = False,
        show_value: bool = True,
    ) -> None:
        """
        Apply data labels to all chart plots.

        Parameters
        ----------
        chart : pptx.chart.chart.Chart
            Chart object to format.
        show_labels : bool, default False
            Whether to show data labels.
        label_position : str, default "center"
            Label position key from LABEL_POSITION.
        number_format : str, default "0.0"
            Data label number format.
        font_size : float, default 5
            Data label font size.
        show_category_name : bool, default False
            Whether to display category names.
        show_value : bool, default True
            Whether to display numeric values.
        """
        validate_option(
            value=label_position,
            options=LABEL_POSITION,
            parameter_name="label_position",
        )

        for plot in chart.plots:
            plot.has_data_labels = show_labels

            if not show_labels:
                continue

            plot.data_labels.show_category_name = show_category_name
            plot.data_labels.show_value = show_value
            plot.data_labels.number_format = number_format
            plot.data_labels.font.name = self.font_name
            plot.data_labels.font.size = Pt(font_size)
            plot.data_labels.position = LABEL_POSITION[label_position]