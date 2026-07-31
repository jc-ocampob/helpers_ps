"""
Utility functions for PowerPoint automation.

This module includes general helpers for color conversion, option validation,
axis scaling, nested list handling and safe text replacement while preserving
PowerPoint formatting where possible.
"""

from __future__ import annotations

from collections.abc import Sequence
from copy import deepcopy
from typing import Any

from pptx.dml.color import RGBColor
from pptx.shapes.group import GroupShape


def hex_to_rgb(hex_color: str) -> RGBColor:
    """
    Convert a hexadecimal color string into a python-pptx RGBColor object.

    Parameters
    ----------
    hex_color : str
        Hexadecimal color string. It can be passed with or without '#'.
        Example: '#122B5F' or '122B5F'.

    Returns
    -------
    RGBColor
        RGBColor object compatible with python-pptx.

    Raises
    ------
    ValueError
        If hex_color is not a valid hexadecimal color.
    """
    if not isinstance(hex_color, str):
        raise ValueError("hex_color must be a string.")

    clean_color = hex_color.strip().lstrip("#")

    if len(clean_color) != 6:
        raise ValueError(
            f"Invalid hex color '{hex_color}'. Expected format is '#RRGGBB'."
        )

    try:
        red = int(clean_color[0:2], 16)
        green = int(clean_color[2:4], 16)
        blue = int(clean_color[4:6], 16)
    except ValueError as exc:
        raise ValueError(
            f"Invalid hex color '{hex_color}'. Expected format is '#RRGGBB'."
        ) from exc

    return RGBColor(red, green, blue)


def validate_option(value: str, options: dict, parameter_name: str) -> None:
    """
    Validate that a value exists in a dictionary of valid options.

    Parameters
    ----------
    value : str
        Value to validate.
    options : dict
        Dictionary where valid values are represented by keys.
    parameter_name : str
        Name of the parameter being validated.

    Raises
    ------
    ValueError
        If value is not a valid key in options.
    """
    if value not in options:
        raise ValueError(
            f"Invalid {parameter_name} '{value}'. "
            f"Available options are: {list(options.keys())}."
        )


def apply_axis_adjustment(
    value: float,
    adjustment_factor: float,
    is_minimum: bool,
) -> float:
    """
    Apply a percentage adjustment to an axis boundary.

    Parameters
    ----------
    value : float
        Axis boundary value.
    adjustment_factor : float
        Percentage adjustment to apply.
    is_minimum : bool
        Whether the value corresponds to the minimum axis boundary.

    Returns
    -------
    float
        Adjusted axis boundary.
    """
    if value == 0:
        return 0.0

    factor = adjustment_factor / 100

    if is_minimum:
        if value < 0:
            return value * (1 + factor)
        return value * (1 - factor)

    if value < 0:
        return value * (1 - factor)

    return value * (1 + factor)


def calculate_axis_bounds(
    values: Sequence[float],
    minimum_adjustment_factor: float = 5.0,
    maximum_adjustment_factor: float = 5.0,
) -> tuple[float, float]:
    """
    Calculate adjusted minimum and maximum axis bounds.

    Parameters
    ----------
    values : Sequence[float]
        Numeric values used to calculate the axis range.
    minimum_adjustment_factor : float, default 5.0
        Percentage adjustment applied to the minimum value.
    maximum_adjustment_factor : float, default 5.0
        Percentage adjustment applied to the maximum value.

    Returns
    -------
    tuple[float, float]
        Adjusted minimum and maximum axis values.

    Raises
    ------
    ValueError
        If values is empty.
    """
    clean_values = [float(value) for value in values if value is not None]

    if not clean_values:
        raise ValueError("values cannot be empty.")

    minimum_value = min(clean_values)
    maximum_value = max(clean_values)

    adjusted_minimum = apply_axis_adjustment(
        value=minimum_value,
        adjustment_factor=minimum_adjustment_factor,
        is_minimum=True,
    )
    adjusted_maximum = apply_axis_adjustment(
        value=maximum_value,
        adjustment_factor=maximum_adjustment_factor,
        is_minimum=False,
    )

    if adjusted_minimum == adjusted_maximum:
        if adjusted_minimum == 0:
            return -1.0, 1.0

        padding = abs(adjusted_minimum) * 0.1
        return adjusted_minimum - padding, adjusted_maximum + padding

    return adjusted_minimum, adjusted_maximum


def flatten_nested_values(values: Sequence[Sequence[float]]) -> list:
    """
    Flatten a nested sequence of numeric values into a single list.

    Parameters
    ----------
    values : Sequence[Sequence[float]]
        Nested numeric sequence.

    Returns
    -------
    list[float]
        Flattened list of values.
    """
    return [item for sublist in values for item in sublist]


def validate_equal_length(
    values: Sequence,
    expected_length: int,
    parameter_name: str,
) -> None:
    """
    Validate that a sequence has an expected length.

    Parameters
    ----------
    values : Sequence
        Sequence to validate.
    expected_length : int
        Expected number of elements.
    parameter_name : str
        Name of the parameter being validated.

    Raises
    ------
    ValueError
        If the sequence length differs from expected_length.
    """
    if len(values) != expected_length:
        raise ValueError(
            f"{parameter_name} must have length {expected_length}. "
            f"Got {len(values)}."
        )


def replace_text_preserve_style(
    target: Any,
    new_text: Any,
    *,
    preserve_linebreaks: bool = True,
    linebreak_replacement: str = " ",
    clear_extra_runs: bool = True,
    clear_extra_paragraphs: bool = True,
    none_as_empty: bool = True,
) -> None:
    """
    Replace text in a PowerPoint shape or table cell while preserving the
    formatting of the first paragraph and first run.

    The function works with any python-pptx object that exposes a ``text_frame``
    attribute, including shapes and table cells.

    Parameters
    ----------
    target : Any
        PowerPoint shape or table cell containing a text frame.
    new_text : Any
        New text to assign. Non-string values are converted to string.
    preserve_linebreaks : bool, default True
        If True, line breaks are preserved by creating one paragraph per line.
        If False, line breaks are replaced by ``linebreak_replacement``.
    linebreak_replacement : str, default " "
        Text used to replace line breaks when ``preserve_linebreaks`` is False.
    clear_extra_runs : bool, default True
        Whether to remove extra runs from the first paragraph.
    clear_extra_paragraphs : bool, default True
        Whether to remove all paragraphs after the first one before adding
        new multiline content.
    none_as_empty : bool, default True
        If True, None is converted to an empty string. If False, None is
        converted to the string ``"None"``.

    Raises
    ------
    ValueError
        If the target does not contain a text frame.
    """
    text_frame = _get_text_frame(target)
    text = _normalize_text(
        new_text,
        preserve_linebreaks=preserve_linebreaks,
        linebreak_replacement=linebreak_replacement,
        none_as_empty=none_as_empty,
    )

    lines = text.split("\n") if preserve_linebreaks else [text]

    if not text_frame.paragraphs:
        text_frame.text = text
        return

    first_paragraph = text_frame.paragraphs[0]
    first_run = _get_or_add_first_run(first_paragraph)

    template_paragraph = first_paragraph
    template_run = first_run

    first_run.text = lines[0] if lines else ""

    if clear_extra_runs:
        _remove_extra_runs(first_paragraph)

    if clear_extra_paragraphs:
        _remove_extra_paragraphs(text_frame)

    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        _copy_paragraph_format(template_paragraph, paragraph)

        run = paragraph.add_run()
        _copy_run_format(template_run, run)
        run.text = line


def replace_shape_text_preserve_style(
    shape: Any,
    new_text: Any,
    *,
    preserve_linebreaks: bool = True,
    linebreak_replacement: str = " ",
    none_as_empty: bool = True,
) -> None:
    """
    Replace text in a PowerPoint shape while preserving existing style.

    Parameters
    ----------
    shape : Any
        PowerPoint shape containing a text frame.
    new_text : Any
        New text to assign.
    preserve_linebreaks : bool, default True
        If True, line breaks are preserved as separate paragraphs.
    linebreak_replacement : str, default " "
        Replacement used when line breaks are not preserved.
    none_as_empty : bool, default True
        If True, None is converted to an empty string.
    """
    replace_text_preserve_style(
        shape,
        new_text,
        preserve_linebreaks=preserve_linebreaks,
        linebreak_replacement=linebreak_replacement,
        none_as_empty=none_as_empty,
    )


def replace_cell_text_preserve_style(
    cell: Any,
    new_text: Any,
    *,
    preserve_linebreaks: bool = True,
    linebreak_replacement: str = " ",
    none_as_empty: bool = True,
) -> None:
    """
    Replace text in a PowerPoint table cell while preserving existing style.

    Parameters
    ----------
    cell : Any
        PowerPoint table cell object.
    new_text : Any
        New text to assign.
    preserve_linebreaks : bool, default True
        If True, line breaks are preserved as separate paragraphs.
    linebreak_replacement : str, default " "
        Replacement used when line breaks are not preserved.
    none_as_empty : bool, default True
        If True, None is converted to an empty string.
    """
    replace_text_preserve_style(
        cell,
        new_text,
        preserve_linebreaks=preserve_linebreaks,
        linebreak_replacement=linebreak_replacement,
        none_as_empty=none_as_empty,
    )


def _get_text_frame(target: Any):
    """
    Return the text frame from a PowerPoint shape or table cell.

    Parameters
    ----------
    target : Any
        Object expected to expose a ``text_frame`` attribute.

    Returns
    -------
    TextFrame
        The PowerPoint text frame.

    Raises
    ------
    ValueError
        If the target does not contain a text frame.
    """
    if hasattr(target, "has_text_frame") and not target.has_text_frame:
        name = getattr(target, "name", "<unnamed>")
        raise ValueError(f"The shape '{name}' does not have a text frame.")

    if not hasattr(target, "text_frame"):
        raise ValueError("The target object does not have a text frame.")

    return target.text_frame


def _normalize_text(
    value: Any,
    *,
    preserve_linebreaks: bool,
    linebreak_replacement: str,
    none_as_empty: bool,
) -> str:
    """
    Normalize input text before assigning it to a PowerPoint text frame.

    Parameters
    ----------
    value : Any
        Raw value to normalize.
    preserve_linebreaks : bool
        Whether line breaks should be preserved.
    linebreak_replacement : str
        Replacement used when line breaks are not preserved.
    none_as_empty : bool
        Whether None should be converted to an empty string.

    Returns
    -------
    str
        Normalized text.
    """
    if value is None and none_as_empty:
        text = ""
    else:
        text = str(value)

    text = text.replace("\r\n", "\n").replace("\r", "\n")

    if not preserve_linebreaks:
        text = text.replace("\n", linebreak_replacement)

    return text


def _get_or_add_first_run(paragraph: Any):
    """
    Return the first run of a paragraph, creating one if necessary.

    Parameters
    ----------
    paragraph : Any
        PowerPoint paragraph object.

    Returns
    -------
    Run
        First run in the paragraph.
    """
    if paragraph.runs:
        return paragraph.runs[0]

    return paragraph.add_run()


def _remove_extra_runs(paragraph: Any) -> None:
    """
    Remove all runs after the first run from a PowerPoint paragraph.

    Parameters
    ----------
    paragraph : Any
        PowerPoint paragraph object.
    """
    for run in list(paragraph.runs[1:]):
        run._r.getparent().remove(run._r)


def _remove_extra_paragraphs(text_frame: Any) -> None:
    """
    Remove all paragraphs after the first paragraph from a text frame.

    Parameters
    ----------
    text_frame : Any
        PowerPoint text frame object.
    """
    while len(text_frame.paragraphs) > 1:
        paragraph = text_frame.paragraphs[-1]
        paragraph._element.getparent().remove(paragraph._element)


def _copy_paragraph_format(source_paragraph: Any, target_paragraph: Any) -> None:
    """
    Copy paragraph-level XML formatting from one paragraph to another.

    This preserves properties such as alignment, indentation, bullet settings,
    paragraph level, spacing, and other paragraph-level styles where available.

    Parameters
    ----------
    source_paragraph : Any
        Paragraph to copy formatting from.
    target_paragraph : Any
        Paragraph to copy formatting to.
    """
    source_ppr = source_paragraph._p.pPr

    if source_ppr is None:
        return

    target_ppr = target_paragraph._p.get_or_add_pPr()
    target_paragraph._p.remove(target_ppr)
    target_paragraph._p.insert(0, deepcopy(source_ppr))


def _copy_run_format(source_run: Any, target_run: Any) -> None:
    """
    Copy run-level XML formatting from one run to another.

    This preserves properties such as font family, size, color, bold, italic,
    underline, and other run-level styles where available.

    Parameters
    ----------
    source_run : Any
        Run to copy formatting from.
    target_run : Any
        Run to copy formatting to.
    """
    source_rpr = source_run._r.rPr

    if source_rpr is None:
        return

    target_rpr = target_run._r.get_or_add_rPr()
    target_run._r.remove(target_rpr)
    target_run._r.insert(0, deepcopy(source_rpr))


def get_shape_by_name(slide, target_name: str):
    """
    Find a shape by name in a slide, including shapes inside grouped shapes.

    Parameters
    ----------
    slide
        PowerPoint slide object.
    target_name : str
        Shape name to search for.

    Returns
    -------
    object | None
        Matching shape if found; otherwise None.
    """

    def walk_shapes(shapes):
        for shape in shapes:
            if shape.name == target_name:
                return shape

            if isinstance(shape, GroupShape):
                found_shape = walk_shapes(shape.shapes)
                if found_shape is not None:
                    return found_shape

        return None

    return walk_shapes(slide.shapes)


def get_layout_by_name(prs, layout_name: str):
    """
    Get a PowerPoint slide layout by name.

    Parameters
    ----------
    prs
        PowerPoint Presentation object.
    layout_name : str
        Layout name to search for.

    Returns
    -------
    pptx.slide.SlideLayout
        Matching slide layout.

    Raises
    ------
    ValueError
        If the layout name is not found.
    """
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout

    raise ValueError(f"Layout '{layout_name}' not found.")

def delete_slide(prs, slide_index):
    slide_id_list = prs.slides._sldIdLst  # underlying XML
    slide_id_list.remove(slide_id_list[slide_index])