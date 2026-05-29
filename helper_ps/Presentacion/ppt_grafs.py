from helper_ps.Presentacion.ppt_dicts import *
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.chart.data import CategoryChartData, BubbleChartData
from pptx.util import Cm, Pt
from pptx.enum.chart import XL_CHART_TYPE
from helper_ps.Presentacion.ppt_helpers import hex_to_rgb
from helper_ps.Config.var_globs import PALETA_COLORES

# Clase para generar grafico de lineas en una ppt
@dataclass
class Line_chart():
    """
    Clase para graficos de lineas en pptx

    Paramentros
    -----------
    presentation: Presentation = None
        Presentación de PowerPoint donde se insertará el gráfico.
    slide: int = 0
        Índice de la diapositiva donde se insertará el gráfico.
    series_names: list[str] = [None]
        Lista de nombres de las series a graficar.
    x: list[list] = [None]
        Lista de listas con los valores del eje x para cada serie.
    y: list[list] = [None]
        Lista de listas con los valores del eje y para cada serie.
    colors: list[str] = PALETA_COLORES_CREDICORP
        Lista de colores en formato hexadecimal para cada serie.
    position_x: int = 5
        Posición horizontal del gráfico en centímetros.
    position_y: int = 5
        Posición vertical del gráfico en centímetros.
    width: int = 10
        Ancho del gráfico en centímetros.
    height: int = 10
        Alto del gráfico en centímetros.
    title_show: bool = True
        Indica si se muestra el título del gráfico.
    title: str = ""
        Texto del título del gráfico.
    title_font: str = "Inter"
        Fuente del título del gráfico.
    title_font_size: int = 6
        Tamaño de la fuente del título del gráfico.
    x_axis_titulo_mostrar: bool = True
        Indica si se muestra el título del eje x.
    x_axis_titulo: str = ""
        Texto del título del eje x.
    x_axis_titulo_fuente: str = "Inter" 
        Fuente del título del eje x.
    x_axis_titulo_tamano: int = 6
        Tamaño de la fuente del título del eje x.
    x_axis_formato: str = "0.00"
        Formato de los valores del eje x.
    y_axis_titulo_mostrar: bool = True
        Indica si se muestra el título del eje y.
    y_axis_titulo: str = ""
        Texto del título del eje y.
    y_axis_titulo_fuente: str = "Inter"
        Fuente del título del eje y.
    y_axis_titulo_tamano: int = 6
        Tamaño de la fuente del título del eje y.
    y_axis_lineas: bool = False
        Indica si se muestran las líneas de la cuadrícula del eje y.
    y_axis_formato: str = "0.00"
        Formato de los valores del eje y.
    legend_show: bool = False
        Indica si se muestra la leyenda del gráfico.
    legend_font_size: int = 5
        Tamaño de la fuente de la leyenda del gráfico.
    legend_position: str = "bottom"
        Posición de la leyenda del gráfico (bottom, top, right, left).

    Funciones
    ---------
    chart_line(self) -> None:
        Función para graficar un gráfico de líneas en una presentación de PowerPoint utilizando los parámetros definidos en la clase.
    """

    presentation: Presentation = None
    slide: int = 0

    series_names: list[str] = field(default_factory=list)
    x: list[list] = field(default_factory=list)
    y: list[list] = field(default_factory=list)
    colors: list[str] = field(default_factory=lambda: PALETA_COLORES)

    position_x: int = 5
    position_y: int = 5

    width: int = 10
    height: int = 10

    title_show: bool = False
    title: str = ""
    title_font: str = "Inter"
    title_font_size: int = 6

    x_axis_titulo_mostrar: bool = False
    x_axis_titulo: str = ""
    x_axis_titulo_tamano: int = 6
    x_axis_formato: str = "0.00"
    x_axis_position: str = "center"

    y_axis_titulo_mostrar: bool = False
    y_axis_titulo: str = ""
    y_axis_titulo_tamano: int = 6
    y_axis_lineas: bool = False
    y_axis_formato: str = "0.00"
    y_factor_ajuste_min: int  = 5.0
    y_factor_ajuste_max: int  = 5.0

    legend_show: bool = False
    legend_font_size: int = 5
    legend_position: str = "bottom"

    # post init para validar alguna opción
    def __post_init__(self):
        if self.legend_position not in LEGEND_POSITION.keys():
            raise TypeError(f"Posición del grafico no es correcta las opciones son {LEGEND_POSITION.keys()}")

    def chart(self) -> None:
        """
        Función para poder graficar un gráfico de líneas en una presentación de PowerPoint utilizando los paramentros de creación definidos en la clase
        """

        # Chart data type
        chart_data = CategoryChartData()

        # save min and max value for y_axis adjustment
        _min_value = []
        _max_value = []

        # For loop to add data to chart
        for serie in range(len(self.series_names)):
            chart_data.categories = self.x[serie]
            chart_data.add_series(self.series_names[serie], self.y[serie])
            _min_value.append(min(self.y[serie]))
            _max_value.append(max(self.y[serie]))
        
        # Definition of y_axis high and low
        _min_value = min(_min_value)
        _max_value = max(_max_value)

        if _min_value < 0:
            _min_value = _min_value * (1 + self.y_factor_ajuste_min/100)
        elif _min_value > 0:
            _min_value = _min_value * (1 - self.y_factor_ajuste_min/100)
        else:
            _min_value = 0
            
        if _max_value < 0:
            _max_value = _max_value * (1 - self.y_factor_ajuste_max/100)
        elif _max_value > 0:
            _max_value = _max_value * (1 + self.y_factor_ajuste_max/100)
        else:
            _max_value = 0

        # Slide where the chart will be put
        slide = self.presentation.slides[self.slide]
        # Create Chort on the presentation
        x = Cm(self.position_x)
        y = Cm(self.position_y)
        w = Cm(self.width)
        h = Cm(self.height)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, chart_data).chart

        # Title
        chart.has_title = self.title_show
        if self.title_show:
            chart.chart_title.text_frame.text = self.title
            chart.chart_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(self.title_font_size)
        
        # Y-axis
        chart.value_axis.has_title = self.y_axis_titulo_mostrar
        if self.y_axis_titulo_mostrar:
            chart.value_axis.axis_title.text_frame.text = self.y_axis_titulo
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.y_axis_titulo_tamano)
            
        chart.value_axis.tick_labels.font.size = Pt(self.y_axis_titulo_tamano)
        chart.value_axis.tick_labels.font.name = self.title_font
        chart.value_axis.minimum_scale = _min_value
        chart.value_axis.has_major_gridlines = self.y_axis_lineas
        chart.value_axis.tick_labels.number_format = self.y_axis_formato

        # X Axis
        chart.category_axis.has_title = self.x_axis_titulo_mostrar
        chart.category_axis.tick_label_position = AXIS_POSITION[self.x_axis_position]
        if self.x_axis_titulo_mostrar:
            chart.category_axis.axis_title.text_frame.text = self.x_axis_titulo
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.x_axis_titulo_tamano)
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font
        chart.category_axis.tick_labels.font.size = Pt(self.x_axis_titulo_tamano)
        chart.category_axis.tick_labels.font.name = self.title_font
        chart.category_axis.tick_labels.number_format = self.x_axis_formato
        
        # Legend
        chart.has_legend = self.legend_show
        if self.legend_show:
            chart.legend.position = LEGEND_POSITION[self.legend_position]
            chart.legend.font.size = Pt(self.legend_font_size)
            chart.legend.font.name = self.title_font

        # Series colors
        _col_c = 0
        for serie in chart.series:
            _s = serie.format.line.color.rgb = hex_to_rgb(self.colors[_col_c])
            _col_c += 1
        
        return 0

# Clase para generar grafico pie en una ppt
@dataclass
class Pie_chart():
    presentation: Presentation = None
    slide: int = 0

    series_name: str = ""
    x: list[str] = field(default_factory=list)
    y: list[int] = field(default_factory=list)

    colors: list[str] = field(default_factory=lambda: PALETA_COLORES)

    position_x: int = 5
    position_y: int = 5

    width: int = 10
    height: int = 10

    title_show: bool = True 
    title: str =  ""
    title_font: str =  "Inter"
    title_font_size: int = 6

    # Legend information
    legend_show: bool = False
    legend_font_size: int = 5
    legend_position: str = "bottom"

    # Data labels
    data_labels_show: bool = False
    data_labels_position: str = "center"
    data_labels_format: str = "0.0"
    data_labels_font_size: int = 5

    # post init para validar alguna opción
    def __post_init__(self) -> None:
        if self.legend_position not in LABEL_POSITION.keys():
            raise TypeError(f"Posición del grafico no es correcta las opciones son {LABEL_POSITION.keys()}")
        return None

    def chart(self):
        # Chart data type
        chart_data = CategoryChartData()
        
        chart_data.categories = self.x
        chart_data.add_series(self.series_name, self.y)

        # Slide where the chart will be put
        slide = self.presentation.slides[self.slide]

        # Create Chort on the presentation
        x = Cm(self.position_x)
        y = Cm(self.position_y)
        w = Cm(self.width)
        h = Cm(self.height)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, chart_data).chart

        # Title
        chart.has_title = self.title_show
        if self.title_show:
            chart.chart_title.text_frame.text = self.title
            chart.chart_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(self.title_font_size)

        # Legend
        chart.has_legend = self.legend_show
        if self.legend_show:
            chart.legend.font.size = Pt(self.legend_font_size)
            chart.legend.position = LEGEND_POSITION[self.legend_position]
        
        # Add data labels
        for plot in chart.plots:
            plot.has_data_labels = self.data_labels_show
            if self.data_labels_show:
                plot.data_labels.show_category_name = False
                plot.data_labels.show_value = True
                plot.data_labels.number_format = self.data_labels_format
                plot.data_labels.font.name = self.title_font
                plot.data_labels.font.size = Pt(self.data_labels_font_size)
                plot.data_labels.position = LABEL_POSITION[self.data_labels_position]

        return None

# Clase para generar grafico de barras
@dataclass
class Bar_chart():
    presentation: Presentation = None
    slide: int = 0

    series_names: list[str] = None
    x: list[str] = None
    y: list[int] = None

    colors: list[str] = None

    position_x: int = 5
    position_y: int = 5

    width: int = 10
    height: int = 10

    title_show: bool = True 
    title: str =  ""
    title_font: str =  "Inter"
    title_font_size: int = 6

    # Legend information
    legend_show: bool = False
    legend_font_size: int = 5
    legend_position: str = "bottom"

    # Data labels
    data_labels_show: bool = False
    data_labels_position: str = "center"
    data_labels_format: str = "0.0"
    data_labels_font_size: int = 5

    x_axis_titulo_mostrar: bool = False
    x_axis_titulo: str = ""
    x_axis_titulo_tamano: int = 6
    x_axis_formato: str = "0.00"
    x_axis_position: str = "center"

    y_axis_titulo_mostrar: bool = False
    y_axis_titulo: str = ""
    y_axis_titulo_tamano: int = 6
    y_axis_lineas: bool = False
    y_axis_formato: str = "0.00"
    y_axis_position: str = "center"
    y_factor_ajuste_min: int  = 10.0
    y_factor_ajuste_max: int  = 10.0

    def chart(self, chart_type: str = "COLUMN_CLUSTERED") -> None:
        # Chart data type
        chart_data = CategoryChartData()

        # save min and max value for y_axis adjustment
        _min_value = []
        _max_value = []
        
        # For loop to add data to chart
        for serie in range(len(self.series_names)):
            chart_data.categories = self.x[serie]
            chart_data.add_series(self.series_names[serie], self.y[serie])
            _min_value.append(min(self.y[serie]))
            _max_value.append(max(self.y[serie]))
        
        # Definition of y_axis high and low
        _min_value = min(_min_value)
        _max_value = max(_max_value)

        if _min_value < 0:
            _min_value = _min_value * (1 + self.y_factor_ajuste_min/100)
        elif _min_value > 0:
            _min_value = _min_value * (1 - self.y_factor_ajuste_min/100)
        else:
            _min_value = 0
            
        if _max_value < 0:
            _max_value = _max_value * (1 - self.y_factor_ajuste_max/100)
        elif _max_value > 0:
            _max_value = _max_value * (1 + self.y_factor_ajuste_max/100)
        else:
            _max_value = 0

        # Slide where the chart will be put
        slide = self.presentation.slides[self.slide]

        # Create Chort on the presentation
        x = Cm(self.position_x)
        y = Cm(self.position_y)
        w = Cm(self.width)
        h = Cm(self.height)

        chart_types = {
            "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
            "COLUMN_STACKED_100": XL_CHART_TYPE.COLUMN_STACKED_100,
            "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
            "BAR_STACKED": XL_CHART_TYPE.BAR_STACKED,
            "BAR_STACKED_100": XL_CHART_TYPE.BAR_STACKED_100,
        }

        if chart_type not in chart_types.keys():
            raise TypeError(f"Tipo de grafico no existe las opciones son {chart_types.keys()}")

        chart = slide.shapes.add_chart(chart_types[chart_type], x, y, w, h, chart_data).chart

        # suponiendo que ya tienes: chart = shape.chart
        for s in chart.series:
            # desactiva "Invertir si es negativo"
            s.invert_if_negative = False

        # Title
        chart.has_title = self.title_show
        if self.title_show:
            chart.chart_title.text_frame.text = self.title
            chart.chart_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(self.title_font_size)

        # Legend
        chart.has_legend = self.legend_show
        if self.legend_show:
            chart.legend.include_in_layout = self.legend_show
            chart.legend.font.size = Pt(self.legend_font_size)
            chart.legend.position = LEGEND_POSITION[self.legend_position]
        
        # Add data labels
        for plot in chart.plots:
            plot.has_data_labels = self.data_labels_show
            if self.data_labels_show:
                plot.data_labels.show_category_name = False
                plot.data_labels.show_value = True
                plot.data_labels.number_format = self.data_labels_format
                plot.data_labels.font.name = self.title_font
                plot.data_labels.font.size = Pt(self.data_labels_font_size)
                plot.data_labels.position = LABEL_POSITION[self.data_labels_position]
        
        # Y-axis
        chart.value_axis.has_title = self.y_axis_titulo_mostrar
        chart.value_axis.tick_label_position = AXIS_POSITION[self.y_axis_position]
        if self.y_axis_titulo_mostrar:
            chart.value_axis.axis_title.text_frame.text = self.y_axis_titulo
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.y_axis_titulo_tamano)
            
        chart.value_axis.tick_labels.font.size = Pt(self.y_axis_titulo_tamano)
        chart.value_axis.tick_labels.font.name = self.title_font
        chart.value_axis.minimum_scale = _min_value
        chart.value_axis.maximum_scale = _max_value
        chart.value_axis.has_major_gridlines = self.y_axis_lineas
        chart.value_axis.tick_labels.number_format = self.y_axis_formato

        # X Axis
        chart.category_axis.has_title = self.x_axis_titulo_mostrar
        chart.category_axis.tick_label_position = AXIS_POSITION[self.x_axis_position]
        if self.x_axis_titulo_mostrar:
            chart.category_axis.axis_title.text_frame.text = self.x_axis_titulo
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.x_axis_titulo_tamano)
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font
        chart.category_axis.tick_labels.font.size = Pt(self.x_axis_titulo_tamano)
        chart.category_axis.tick_labels.font.name = self.title_font
        chart.category_axis.tick_labels.number_format = self.x_axis_formato

        return None

# Bubble chart 
@dataclass
class Bubble_chart():

    presentation: Presentation = None
    slide: int = 0

    series_names: list[str] = field(default_factory=list)
    x: list[list] = field(default_factory=list)
    y: list[list] = field(default_factory=list)
    z: list[list] = field(default_factory=list)
    colors: list[str] = field(default_factory=lambda: PALETA_COLORES)

    position_x: int = 5
    position_y: int = 5

    width: int = 10
    height: int = 10

    title_show: bool = False
    title: str = ""
    title_font: str = "Inter"
    title_font_size: int = 6

    x_axis_titulo_mostrar: bool = False
    x_axis_titulo: str = ""
    x_axis_titulo_tamano: int = 6
    x_axis_formato: str = "0.00"
    x_factor_ajuste_min: int = 5.0
    x_factor_ajuste_max: int = 5.0

    y_axis_titulo_mostrar: bool = False
    y_axis_titulo: str = ""
    y_axis_titulo_tamano: int = 6
    y_axis_lineas: bool = False
    y_axis_formato: str = "0.00"
    y_factor_ajuste_min: int  = 5.0
    y_factor_ajuste_max: int  = 5.0

    legend_show: bool = False
    legend_font_size: int = 5
    legend_position: str = "bottom"

    # post init para validar alguna opción
    def __post_init__(self):
        if self.legend_position not in LEGEND_POSITION.keys():
            raise TypeError(f"Posición del grafico no es correcta las opciones son {LEGEND_POSITION.keys()}")

    def chart(self) -> None:
        """
        Función para poder graficar un gráfico de burbujas en una presentación de PowerPoint
        utilizando los parámetros de creación definidos en la clase.
        """

        # Chart data type
        chart_data = BubbleChartData()

        # save min and max value for axis adjustment
        _min_value_y = []
        _max_value_y = []
        _min_value_x = []
        _max_value_x = []

        # For loop to add data to chart
        for serie in range(len(self.series_names)):

            x_vals = self.x[serie]
            y_vals = self.y[serie]  # viene como 5 para 5%
            z_vals = self.z[serie]  # tamaño burbuja (ya viene en escala correcta)

            # Validación básica de longitudes
            if not (len(x_vals) == len(y_vals) == len(z_vals)):
                raise ValueError(
                    f"Longitudes no coinciden en serie '{self.series_names[serie]}': "
                    f"len(x)={len(x_vals)}, len(y)={len(y_vals)}, len(z)={len(z_vals)}"
                )

            # Track min/max for axis scaling
            _min_value_x.append(min(x_vals))
            _max_value_x.append(max(x_vals))
            _min_value_y.append(min(y_vals))
            _max_value_y.append(max(y_vals))

            # BubbleChartData: add_series SOLO recibe el nombre
            s = chart_data.add_series(self.series_names[serie])

            # Para bubble chart, agregas puntos (x, y, size) uno por uno
            for x, y, z in zip(x_vals, y_vals, z_vals):
                s.add_data_point(x, y, z)

        # Definition of y_axis high and low
        _min_value_y = min(_min_value_y)
        _max_value_y = max(_max_value_y)

        if _min_value_y < 0:
            _min_value_y = _min_value_y * (1 + self.y_factor_ajuste_min/100)
        elif _min_value_y > 0:
            _min_value_y = _min_value_y * (1 - self.y_factor_ajuste_min/100)
        else:
            _min_value_y = 0

        if _max_value_y < 0:
            _max_value_y = _max_value_y * (1 - self.y_factor_ajuste_max/100)
        elif _max_value_y > 0:
            _max_value_y = _max_value_y * (1 + self.y_factor_ajuste_max/100)
        else:
            _max_value_y = 0

        # Definition of x_axis high and low
        _min_value_x = min(_min_value_x)
        _max_value_x = max(_max_value_x)

        if _min_value_x < 0:
            _min_value_x = _min_value_x * (1 + self.x_factor_ajuste_min/100)
        elif _min_value_x > 0:
            _min_value_x = _min_value_x * (1 - self.x_factor_ajuste_min/100)
        else:
            _min_value_x = 0

        if _max_value_x < 0:
            _max_value_x = _max_value_x * (1 - self.x_factor_ajuste_max/100)
        elif _max_value_x > 0:
            _max_value_x = _max_value_x * (1 + self.x_factor_ajuste_max/100)
        else:
            _max_value_x = 0

        # Slide where the chart will be put
        slide = self.presentation.slides[self.slide]

        # Create Chart on the presentation
        x = Cm(self.position_x)
        y = Cm(self.position_y)
        w = Cm(self.width)
        h = Cm(self.height)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE, x, y, w, h, chart_data
        ).chart

        # Title
        chart.has_title = self.title_show
        if self.title_show:
            chart.chart_title.text_frame.text = self.title
            chart.chart_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(self.title_font_size)

        # Y-axis
        chart.value_axis.has_title = self.y_axis_titulo_mostrar
        if self.y_axis_titulo_mostrar:
            chart.value_axis.axis_title.text_frame.text = self.y_axis_titulo
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.y_axis_titulo_tamano)

        chart.value_axis.tick_labels.font.size = Pt(self.y_axis_titulo_tamano)
        chart.value_axis.tick_labels.font.name = self.title_font
        chart.value_axis.minimum_scale = float(_min_value_y)
        chart.value_axis.maximum_scale = float(_max_value_y)
        chart.value_axis.has_major_gridlines = self.y_axis_lineas
        chart.value_axis.tick_labels.number_format = self.y_axis_formato

        # X Axis (en bubble chart python-pptx lo expone como category_axis)
        chart.category_axis.has_title = self.x_axis_titulo_mostrar
        if self.x_axis_titulo_mostrar:
            chart.category_axis.axis_title.text_frame.text = self.x_axis_titulo
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(self.x_axis_titulo_tamano)
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = self.title_font

        chart.category_axis.tick_labels.font.size = Pt(self.x_axis_titulo_tamano)
        chart.category_axis.tick_labels.font.name = self.title_font
        chart.category_axis.tick_labels.number_format = self.x_axis_formato
        chart.category_axis.minimum_scale = float(_min_value_x)
        chart.category_axis.maximum_scale = float(_max_value_x)

        # Legend
        chart.has_legend = self.legend_show
        if self.legend_show:
            chart.legend.position = LEGEND_POSITION[self.legend_position]
            chart.legend.font.size = Pt(self.legend_font_size)
            chart.legend.font.name = self.title_font

        # Series colors (relleno sólido para burbujas)
        _col_c = 0
        for serie in chart.series:
            serie.format.fill.solid()
            serie.format.fill.fore_color.rgb = hex_to_rgb(self.colors[_col_c])
            serie.format.line.fill.background()  # sin borde (opcional)
            _col_c += 1

        return 0
