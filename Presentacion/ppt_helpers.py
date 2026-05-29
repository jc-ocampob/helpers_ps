from pptx.dml.color import RGBColor
from copy import deepcopy
from pptx.shapes.group import GroupShape

# ------------------------------------
# Section 1: funciones
# ------------------------------------
# HEX to RGB formato
def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16))

# ------------------------------------
# Section 3: Tablas PPTX
# ------------------------------------
def set_cell_text_preserve_format(cell, new_text):
    """
    Reemplaza el texto visible de una celda intentando preservar el formato existente:
    - Reutiliza el primer run del primer párrafo (mantiene font, size, bold, color, etc.)
    - Elimina texto de runs/párrafos adicionales para no mezclar contenido viejo.
    """
    tf = cell.text_frame

    # Si no hay párrafos (raro), crea uno
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    # Si no hay runs, crea uno
    if len(p.runs) == 0:
        run = p.add_run()
        run.text = new_text
    else:
        # 1) Setear el texto en el primer run (conserva formato)
        p.runs[0].text = new_text

        # 2) Vaciar texto de runs adicionales del primer párrafo
        for r in p.runs[1:]:
            r.text = ""

    # 3) Vaciar párrafos adicionales (si existieran) para evitar texto residual
    for extra_p in tf.paragraphs[1:]:
        for r in extra_p.runs:
            r.text = ""

# ------------------------------------
# Section 3: Cajas de texto PPTX
# ------------------------------------
def set_text_keep_format(shape, new_text: str, keep_linebreaks: bool = True):
    """
    Reemplaza el texto manteniendo el formato original del shape.
    - NO usa text_frame.text (porque resetea estilos con frecuencia).
    - Reutiliza el primer run existente (que ya tiene el estilo correcto).
    - Si new_text tiene saltos de línea y keep_linebreaks=True,
      crea párrafos adicionales copiando el formato del primer párrafo/run.
    """
    if not shape.has_text_frame:
        raise ValueError(f"El shape '{shape.name}' no tiene text_frame.")

    tf = shape.text_frame
    if len(tf.paragraphs) == 0:
        p0 = tf.paragraphs[0]
    else:
        p0 = tf.paragraphs[0]

    # Asegurar que exista al menos un run
    if len(p0.runs) == 0:
        r0 = p0.add_run()
    else:
        r0 = p0.runs[0]

    # Guardar "plantilla" de formato del párrafo y run (XML)
    p0_xml = deepcopy(p0._p)       # formato del párrafo
    r0_xml = deepcopy(r0._r)       # formato del run

    # Limpieza controlada: vaciar texto de todos los runs existentes (pero no borrar runs)
    for r in p0.runs:
        r.text = ""

    # Si no quieres respetar saltos, simple:
    if not keep_linebreaks or ("\n" not in new_text):
        r0.text = new_text
        return

    # Si hay saltos de línea, creamos varios párrafos con mismo estilo
    lines = new_text.split("\n")

    # 1ra línea en el run original
    r0.text = lines[0]

    # Eliminar párrafos extra existentes (si tu placeholder tenía varios)
    # (python-pptx no tiene remove directo; estrategia: setear texto vacío en esos párrafos)
    for p in tf.paragraphs[1:]:
        # vaciamos el texto sin tocar el resto del shape
        for r in p.runs:
            r.text = ""

    # Crear párrafos adicionales copiando el estilo base
    for line in lines[1:]:
        p = tf.add_paragraph()
        # Copiar formato del párrafo base
        p._p[:] = deepcopy(p0_xml[:])
        # Crear run y copiar formato run base
        rr = p.add_run()
        rr._r[:] = deepcopy(r0_xml[:])
        rr.text = line

def find_shape_by_name(slide, target_name: str):
    def _walk(shapes):
        for s in shapes:
            if s.name == target_name:
                return s
            if isinstance(s, GroupShape):
                found = _walk(s.shapes)
                if found is not None:
                    return found
        return None

    shp = _walk(slide.shapes)
    if shp is None:
        raise KeyError(f"No se encontró el shape con nombre: {target_name}")
    return shp


def get_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    raise ValueError(f"Layout '{layout_name}' not found")
